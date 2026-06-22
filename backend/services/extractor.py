"""
Document text extractor — aligned with Data_Wrangling_(NO_OCR_FIX).ipynb
Supports: .docx, .pptx, .pdf, .doc, .ppt
"""
import os
import io
import subprocess
import tempfile
from pathlib import Path


def extract_text(file_bytes: bytes, filename: str) -> str:
    ext = Path(filename).suffix.lower()
    
    # Handle legacy formats via LibreOffice conversion if possible
    if ext in (".doc", ".ppt"):
        return _extract_legacy_via_conversion(file_bytes, ext)
        
    if ext == ".docx":
        return _extract_docx(file_bytes)
    if ext == ".pptx":
        return _extract_pptx(file_bytes)
    if ext == ".pdf":
        return _extract_pdf(file_bytes)
    
    raise ValueError(f"Format tidak didukung: {ext}. Gunakan .docx, .pptx, atau .pdf.")


def _extract_legacy_via_conversion(file_bytes: bytes, ext: str) -> str:
    """Converts .doc/.ppt to .docx/.pptx using LibreOffice (soffice) then extracts."""
    target_format = "docx" if ext == ".doc" else "pptx"
    
    with tempfile.TemporaryDirectory() as tmpdir:
        input_path = Path(tmpdir) / f"temp_input{ext}"
        input_path.write_bytes(file_bytes)
        
        try:
            # Check for soffice path
            import shutil
            soffice_cmd = shutil.which("soffice") or shutil.which("libreoffice") or "soffice"
            
            if os.name == 'nt' and soffice_cmd in ("soffice", "libreoffice"):
                # Common Windows paths if not in PATH
                for path in [r"C:\Program Files\LibreOffice\program\soffice.exe", r"C:\Program Files (x86)\LibreOffice\program\soffice.exe"]:
                    if os.path.exists(path):
                        soffice_cmd = path
                        break

            subprocess.run([
                soffice_cmd, "--headless", "--convert-to", target_format,
                str(input_path), "--outdir", tmpdir
            ], check=True, capture_output=True)
            
            output_path = Path(tmpdir) / f"temp_input.{target_format}"
            if output_path.exists():
                new_bytes = output_path.read_bytes()
                return _extract_docx(new_bytes) if target_format == "docx" else _extract_pptx(new_bytes)
        except Exception as e:
            print(f"LibreOffice conversion failed: {e}")
            
    # Fallback if LibreOffice is not installed or fails
    raise ValueError(
        f"Gagal memproses file {ext}. Pastikan LibreOffice terinstal di sistem/Hugging Face space "
        f"atau konversikan file ke .{target_format} secara manual."
    )


# ── DOCX ─────────────────────────────────────────────────────────

def _extract_docx(file_bytes: bytes) -> str:
    from docx import Document
    doc = Document(io.BytesIO(file_bytes))
    text_content = []
    
    try:
        # Paragraphs with markers
        for para in doc.paragraphs:
            text = para.text.strip()
            if not text: continue
            
            # Logic from notebook: Detect Heading / List
            if para.style.name.startswith('Heading'):
                text_content.append(f"[JUDUL/SUB] {text}")
            elif 'List' in para.style.name:
                text_content.append(f"[NUMBERING] {text}")
            else:
                text_content.append(text)

        # Tables with markers
        for table in doc.tables:
            text_content.append("\n[TABEL START]")
            for row in table.rows:
                row_data = [cell.text.replace('\n', ' ').strip() for cell in row.cells]
                text_content.append(" | ".join(row_data))
            text_content.append("[TABEL END]\n")

        return "\n".join(text_content)
    except Exception as e:
        return f"Error DOCX Extraction: {e}"


# ── PPTX ─────────────────────────────────────────────────────────

def _extract_pptx(file_bytes: bytes) -> str:
    from pptx import Presentation
    text_content = []
    
    try:
        prs = Presentation(io.BytesIO(file_bytes))
        for i, slide in enumerate(prs.slides):
            text_content.append(f"\n--- Slide {i+1} ---")
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text.strip()
                        if text:
                            text_content.append(text)

                if shape.has_table:
                    text_content.append("[TABEL PPTX]")
                    for row in shape.table.rows:
                        row_data = [cell.text_frame.text.replace('\n', ' ').strip() for cell in row.cells]
                        text_content.append(" | ".join(row_data))

        return "\n".join(text_content)
    except Exception as e:
        return f"Error PPTX Extraction: {e}"


# ── PDF ──────────────────────────────────────────────────────────

def _extract_pdf(file_bytes: bytes) -> str:
    # Note: Notebook uses PyMuPDF (fitz), but project currently uses pypdf.
    # Aligning with notebook's logic of extracting text + tables.
    import fitz # PyMuPDF
    text_content = []
    
    try:
        doc = fitz.open(stream=file_bytes, filetype="pdf")
        for page in doc:
            text_content.append(page.get_text("text"))

            # Table extraction (aligned with notebook)
            try:
                tabs = page.find_tables()
                for tab in tabs:
                    text_content.append("\n[TABEL PDF]")
                    for row in tab.extract():
                        row_data = [str(cell).replace('\n', ' ').strip() if cell else "" for cell in row]
                        text_content.append(" | ".join(row_data))
            except:
                pass 

        return "\n".join(text_content)
    except Exception as e:
        return f"Error PDF Extraction: {e}"
